"""Genera el pitch deck de MichiCheck en el formato oficial de la hackatón.

Ocho diapositivas, 5 minutos, con las figuras embebidas y el guion en las
notas del ponente (con su cronómetro por bloque).

El estilo replica la guía del INSN: fondo azul-violeta profundo, acentos de
neón, tipografía blanca y numeración de sección. Las figuras van sobre
tarjetas blancas redondeadas porque se generaron con fondo claro; ponerlas
directamente sobre el fondo oscuro dejaría el texto de los ejes ilegible.

    python scripts/make_ppt.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

RAIZ = Path(__file__).resolve().parents[1]
FIG = RAIZ / "docs" / "figuras"
CAP = RAIZ / "docs" / "capturas"
SALIDA = RAIZ / "MichiCheck_pitch.pptx"

FONDO = RGBColor(0x0B, 0x0A, 0x2E)
FONDO2 = RGBColor(0x2A, 0x0F, 0x4D)
BLANCO = RGBColor(0xFF, 0xFF, 0xFF)
CIAN = RGBColor(0x22, 0xD3, 0xEE)
MAGENTA = RGBColor(0xFF, 0x2D, 0x8F)
VIOLETA = RGBColor(0xA8, 0x55, 0xF7)
NARANJA = RGBColor(0xFF, 0x8C, 0x42)
VERDE = RGBColor(0x22, 0xC5, 0x5E)
GRIS = RGBColor(0x9C, 0xA3, 0xC4)
ROJO = RGBColor(0xEF, 0x44, 0x44)

ANCHO, ALTO = Inches(13.333), Inches(7.5)
ACENTOS = [MAGENTA, CIAN, VIOLETA, NARANJA]


def _fondo(slide) -> None:
    """Fondo degradado + los dos adornos de esquina de la guía."""
    fondo = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, ANCHO, ALTO)
    fondo.line.fill.background()
    relleno = fondo.fill
    relleno.gradient()
    relleno.gradient_angle = 45.0
    relleno.gradient_stops[0].color.rgb = FONDO
    relleno.gradient_stops[1].color.rgb = FONDO2
    fondo.shadow.inherit = False

    arco = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.6), Inches(6.0),
                                  Inches(3.2), Inches(3.2))
    arco.fill.background()
    arco.line.color.rgb = MAGENTA
    arco.line.width = Pt(2.5)
    arco.shadow.inherit = False


def _texto(slide, x, y, w, h, texto, tam=18, color=BLANCO, negrita=False,
           alinear=PP_ALIGN.LEFT, interlineado=1.15, cursiva=False):
    caja = slide.shapes.add_textbox(x, y, w, h)
    marco = caja.text_frame
    marco.word_wrap = True
    marco.vertical_anchor = MSO_ANCHOR.TOP
    for i, linea in enumerate(str(texto).split("\n")):
        p = marco.paragraphs[0] if i == 0 else marco.add_paragraph()
        p.alignment = alinear
        p.line_spacing = interlineado
        r = p.add_run()
        r.text = linea
        r.font.size = Pt(tam)
        r.font.bold = negrita
        r.font.italic = cursiva
        r.font.color.rgb = color
        r.font.name = "Segoe UI"
    return caja


def _encabezado(slide, numero: str, titulo: str, acento=CIAN) -> None:
    """Numeración de sección + título, al estilo de la guía."""
    _texto(slide, Inches(.75), Inches(.42), Inches(1.3), Inches(.9),
           numero, tam=40, color=acento, negrita=True)
    _texto(slide, Inches(1.62), Inches(.45), Inches(10.6), Inches(.9),
           titulo, tam=34, negrita=True)
    linea = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(.78),
                                   Inches(1.28), Inches(1.05), Pt(3.5))
    linea.fill.solid()
    linea.fill.fore_color.rgb = acento
    linea.line.fill.background()
    linea.shadow.inherit = False
    _texto(slide, Inches(9.0), Inches(.3), Inches(3.6), Inches(.5),
           "CONECTANDO IDEAS PARA\nINNOVAR EN LA SALUD INFANTIL",
           tam=8, color=GRIS, alinear=PP_ALIGN.RIGHT, interlineado=1.3)


def _tarjeta(slide, x, y, w, h, borde=CIAN, relleno=None):
    """Tarjeta de bordes redondeados con contorno de neón."""
    f = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    f.adjustments[0] = 0.06
    if relleno is None:
        f.fill.solid()
        f.fill.fore_color.rgb = RGBColor(0x16, 0x14, 0x40)
        f.fill.transparency = 0.35
    else:
        f.fill.solid()
        f.fill.fore_color.rgb = relleno
    f.line.color.rgb = borde
    f.line.width = Pt(1.6)
    f.shadow.inherit = False
    return f


def _tarjeta_texto(slide, x, y, w, h, titulo, cuerpo, borde=CIAN,
                   tam_t=15, tam_c=12):
    _tarjeta(slide, x, y, w, h, borde)
    _texto(slide, x + Inches(.24), y + Inches(.16), w - Inches(.48),
           Inches(.42), titulo, tam=tam_t, color=borde, negrita=True)
    _texto(slide, x + Inches(.24), y + Inches(.62), w - Inches(.48),
           h - Inches(.78), cuerpo, tam=tam_c, color=BLANCO, interlineado=1.22)


def _figura(slide, ruta: Path, x, y, w, margen=Inches(.12), h_max=None,
            centrar=True):
    """Figura sobre panel blanco. Devuelve el alto ocupado.

    Si ``h_max`` limita el alto, la imagen se reescala manteniendo proporción
    en vez de desbordar la diapositiva — que es lo que pasaba con las figuras
    verticales al fijar solo el ancho.
    """
    from PIL import Image

    with Image.open(ruta) as im:
        proporcion = im.height / im.width
    h = Emu(int(w * proporcion))
    if h_max is not None and h > h_max:
        h = Emu(int(h_max))
        nuevo_w = Emu(int(h / proporcion))
        if centrar:
            x = Emu(int(x + (w - nuevo_w) / 2))
        w = nuevo_w
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   x - margen, y - margen,
                                   w + 2 * margen, h + 2 * margen)
    panel.adjustments[0] = 0.03
    panel.fill.solid()
    panel.fill.fore_color.rgb = BLANCO
    panel.line.fill.background()
    panel.shadow.inherit = False
    slide.shapes.add_picture(str(ruta), x, y, width=w)
    return h


def _numero_grande(slide, x, y, w, valor, etiqueta, color=CIAN, tam=54):
    _texto(slide, x, y, w, Inches(1.0), valor, tam=tam, color=color,
           negrita=True, alinear=PP_ALIGN.CENTER)
    _texto(slide, x, y + Inches(.92), w, Inches(.6), etiqueta, tam=12,
           color=GRIS, alinear=PP_ALIGN.CENTER, interlineado=1.15)


def _notas(slide, texto: str) -> None:
    slide.notes_slide.notes_text_frame.text = texto.strip()


def _pie(slide, texto: str, color=GRIS) -> None:
    _texto(slide, Inches(.75), Inches(6.85), Inches(11.8), Inches(.5),
           texto, tam=11, color=color, alinear=PP_ALIGN.CENTER, cursiva=True)


def construir() -> Presentation:
    prs = Presentation()
    prs.slide_width, prs.slide_height = ANCHO, ALTO
    blanco = prs.slide_layouts[6]

    s = prs.slides.add_slide(blanco); _fondo(s)
    _texto(s, Inches(.9), Inches(.45), Inches(6), Inches(.5),
           "HACKATÓN NIÑO SAN BORJA 2026", tam=13, color=CIAN, negrita=True)
    _texto(s, Inches(.9), Inches(1.35), Inches(8.6), Inches(1.6),
           "MichiCheck", tam=76, negrita=True)
    _texto(s, Inches(.9), Inches(2.85), Inches(8.8), Inches(1.0),
           "Que ningún niño abandone su tratamiento\npor la distancia",
           tam=25, color=VIOLETA, interlineado=1.25)

    for i, (etiqueta, valor, color) in enumerate([
        ("Equipo", "[nombre del equipo]", MAGENTA),
        ("Solución", "MichiCheck", CIAN),
        ("Desafío", "3 · Ruta Hematológica", VIOLETA),
        ("Misión", "Continuidad del cuidado\nsin salir de la región", NARANJA),
    ]):
        x = Inches(.9 + i * 2.95)
        _tarjeta(s, x, Inches(4.55), Inches(2.72), Inches(1.55), color)
        _texto(s, x + Inches(.22), Inches(4.72), Inches(2.3), Inches(.35),
               etiqueta.upper(), tam=10, color=color, negrita=True)
        _texto(s, x + Inches(.22), Inches(5.12), Inches(2.32), Inches(.9),
               valor, tam=13, negrita=True, interlineado=1.2)

    _pie(s, "Instituto Nacional de Salud del Niño – San Borja · 2026")
    _notas(s, """
[0:00–0:10] PORTADA

«Somos [equipo] y traemos MichiCheck, para el desafío de Ruta Hematológica.»

No leer las tarjetas. Pasar rápido.
""")

    s = prs.slides.add_slide(blanco); _fondo(s)
    _encabezado(s, "2.", "Gancho", MAGENTA)

    _tarjeta(s, Inches(.75), Inches(1.6), Inches(6.5), Inches(3.5), MAGENTA)
    _texto(s, Inches(1.05), Inches(1.85), Inches(5.9), Inches(3.1),
           "Una niña de 6 años con leucemia vive en Bagua.\n\n"
           "Está en el día 10 después de su quimioterapia — el día "
           "en que sus defensas tocan fondo.\n\n"
           "Tiene 38.6 de fiebre.",
           tam=20, interlineado=1.35)

    _tarjeta(s, Inches(7.6), Inches(1.6), Inches(4.95), Inches(3.5), CIAN)
    _texto(s, Inches(7.95), Inches(1.9), Inches(4.3), Inches(.4),
           "EL DATO", tam=11, color=CIAN, negrita=True)
    _texto(s, Inches(7.95), Inches(2.35), Inches(4.3), Inches(1.2),
           "56.22 %", tam=62, color=CIAN, negrita=True)
    _texto(s, Inches(7.95), Inches(3.6), Inches(4.3), Inches(1.3),
           "de los niños que fallecieron en el instituto\n"
           "NO era de Lima ni Callao\n\nSala Situacional INSNSB, nov. 2025",
           tam=13, color=GRIS, interlineado=1.3)

    _tarjeta(s, Inches(.75), Inches(5.35), Inches(11.8), Inches(1.25), NARANJA)
    _texto(s, Inches(1.1), Inches(5.62), Inches(11.1), Inches(.8),
           "«No reemplazamos el hemograma. Llenamos los veinte días "
           "en que ese niño no tiene ninguno.»",
           tam=21, negrita=True, alinear=PP_ALIGN.CENTER)
    _notas(s, """
[0:10–0:55] GANCHO — 45 s

HISTORIA (20 s):
«Una niña de 6 años con leucemia vive en Bagua. Está en el día 10 después de su
quimioterapia — el día en que sus defensas tocan fondo. Tiene 38.6 de fiebre.
Si su recuento está por debajo de 500, esto es una emergencia oncológica:
necesita antibiótico en la primera hora. Si está por encima, puede esperar.
NADIE EN BAGUA PUEDE SABER CUÁL DE LAS DOS COSAS ES.»

EL DATO (10 s):
«De los niños que fallecieron en el instituto, el 56% no era de Lima. Es el dato
de su propia Sala Situacional.»

POR QUÉ AHORA (8 s):
«Porque el instituto ya tiene un Plan de Leucemias aprobado cuya meta es atender
al 50% de estos pacientes EN SUS REGIONES. La decisión ya está tomada. Falta el
instrumento.»

FRASE (7 s): leer la de abajo, despacio.
""")

    s = prs.slides.add_slide(blanco); _fondo(s)
    _encabezado(s, "3.", "Problema e impacto", MAGENTA)
    _figura(s, FIG / "07_carga_de_viajes.png", Inches(.95), Inches(1.62),
            Inches(7.5), h_max=Inches(4.9))

    _tarjeta_texto(s, Inches(8.85), Inches(1.5), Inches(3.7), Inches(1.62),
                   "¿CÓMO SE ABORDA HOY?",
                   "El INSN ya tiene Comité de Abandono y usa la app IMPACTO.\n"
                   "Abandono nacional: 18.6 % → 8.5 %",
                   CIAN, tam_t=12, tam_c=12)
    _tarjeta_texto(s, Inches(8.85), Inches(3.28), Inches(3.7), Inches(1.82),
                   "¿QUÉ FALTA?",
                   "IMPACTO llama DESPUÉS de que el niño faltó.\n\n"
                   "El 8.5 % que queda no falta por olvido: seguro público, "
                   "ruralidad y distancia son estructurales.",
                   MAGENTA, tam_t=12, tam_c=11.5)
    _tarjeta_texto(s, Inches(8.85), Inches(5.26), Inches(3.7), Inches(1.4),
                   "EVIDENCIA",
                   "Hematología = 47.21 % del cáncer infantil del INSNSB\n"
                   "LLA = 1ª causa de muerte (15.21 %)",
                   NARANJA, tam_t=12, tam_c=11.5)
    _notas(s, """
[0:55–1:55] PROBLEMA E IMPACTO — 60 s   ← gana el 25 % de la rúbrica

«Antes de proponer nada: el INSN YA TIENE un Comité de Abandono y usa la
aplicación IMPACTO. Y funcionó — el abandono nacional bajó a la mitad. Llamar al
que faltó recupera al que se olvidó.

Pero el 8.5 % que queda NO FALTA POR OLVIDO. Los tres predictores del abandono
son estructurales: seguro público, ruralidad, vivir fuera de la capital. Ninguna
llamada telefónica cambia ninguno de los tres.

Esa niña hará 45 VIAJES EN 18 MESES. Son 810 horas y el 54 % DEL INGRESO DE SU
FAMILIA. El abandono no es una decisión: es una acumulación hasta que ya no se
puede.

Nosotros no llamamos al que faltó. HACEMOS QUE TENGA QUE VENIR MENOS VECES.»
""")

    s = prs.slides.add_slide(blanco); _fondo(s)
    _encabezado(s, "4.", "Solución y MVP", VIOLETA)
    _texto(s, Inches(.78), Inches(1.42), Inches(11.8), Inches(.5),
           "Tamizaje óptico no invasivo que permite controlar el tratamiento "
           "en la posta, sin aguja ni laboratorio",
           tam=17, color=VIOLETA, negrita=True)
    _figura(s, FIG / "01_umbral_pediatrico.png", Inches(.95), Inches(2.15),
            Inches(7.3), h_max=Inches(4.3))

    _tarjeta(s, Inches(8.7), Inches(2.0), Inches(3.85), Inches(3.6), MAGENTA)
    _texto(s, Inches(9.0), Inches(2.2), Inches(3.3), Inches(.4),
           "¿QUÉ LA HACE DIFERENTE?", tam=12, color=MAGENTA, negrita=True)
    _texto(s, Inches(9.0), Inches(2.72), Inches(3.3), Inches(2.7),
           "El método cuenta leucocitos TOTALES.\n\n"
           "A los 2 años solo el 31 % son neutrófilos (59 % en el adulto).\n\n"
           "Con el umbral adulto de 7 gaps/min la alerta salta en:",
           tam=12.5, interlineado=1.3)
    _texto(s, Inches(9.0), Inches(4.62), Inches(1.6), Inches(.9),
           "487", tam=34, color=VERDE, negrita=True, alinear=PP_ALIGN.CENTER)
    _texto(s, Inches(10.75), Inches(4.62), Inches(1.6), Inches(.9),
           "272", tam=34, color=ROJO, negrita=True, alinear=PP_ALIGN.CENTER)
    _texto(s, Inches(9.0), Inches(5.16), Inches(1.6), Inches(.4),
           "adulto ✓", tam=11, color=GRIS, alinear=PP_ALIGN.CENTER)
    _texto(s, Inches(10.75), Inches(5.16), Inches(1.6), Inches(.4),
           "niño 2 años ✗", tam=11, color=GRIS, alinear=PP_ALIGN.CENTER)

    _pie(s, "Se pierde entera la franja 272–500, que es donde hay que actuar",
         NARANJA)
    _notas(s, """
[1:55–3:10] SOLUCIÓN Y MVP — 75 s   ← el momento fuerte

CÓMO FUNCIONA (25 s):
«A 530 nanómetros la hemoglobina absorbe luz: los glóbulos rojos se ven negros.
Un glóbulo blanco no tiene hemoglobina — deja pasar la luz y deja un HUECO
BRILLANTE que viaja por el capilar. Contar huecos es contar glóbulos blancos.
Y la relación no es aprendida: es geometría. Con los parámetros del paper de
referencia nuestro modelo predice 3.773 células. El paper reporta 3.773. NO
AJUSTAMOS NADA.»

QUÉ LA HACE DIFERENTE (35 s):
«Esto ya existe para adultos: PointCheck, con designación FDA Breakthrough.
Pero cuenta glóbulos blancos TOTALES, y en un niño de 2 años solo el 31 % son
neutrófilos frente al 59 % del adulto.
El umbral comercial de 7 huecos por minuto equivale, en un adulto, a un recuento
de 487: correcto. EN UN NIÑO DE 2 AÑOS ESE MISMO UMBRAL RECIÉN SE DISPARA CUANDO
EL RECUENTO YA CAYÓ A 272.»

>>> PAUSA DE DOS SEGUNDOS AQUÍ <<<

SEGUNDA SEÑAL (15 s):
«Y del mismo vídeo sale otra cosa: en mantenimiento el objetivo ES tener el
recuento entre 500 y 1500. Un resultado alto NO es buena noticia. El 44 % de
estos niños tiene adherencia insuficiente y eso TRIPLICA el riesgo de recaída.
IMPACTO registra si el niño vino a la cita; no si tomó la pastilla.»

SI VAS TARDE: recortar este último bloque.
""")

    s = prs.slides.add_slide(blanco); _fondo(s)
    _encabezado(s, "5.", "Viabilidad técnica y económica", CIAN)

    for i, (valor, etiqueta, color) in enumerate([
        ("S/ 226–401", "por unidad\n(citómetro: decenas de miles USD)", CIAN),
        ("0", "reactivos y consumibles", VERDE),
        ("Sin internet", "funciona offline\nen la posta", VIOLETA),
        ("Clase I", "DIGEMID\nregistro simplificado", NARANJA),
    ]):
        x = Inches(.8 + i * 3.05)
        _tarjeta(s, x, Inches(1.55), Inches(2.82), Inches(2.05), color)
        _numero_grande(s, x, Inches(1.78), Inches(2.82), valor, etiqueta,
                       color, tam=30 if len(valor) > 6 else 44)

    _tarjeta_texto(
        s, Inches(.8), Inches(3.85), Inches(5.95), Inches(2.75),
        "RUTA REGULATORIA",
        "La Ley 29459 incluye al SOFTWARE en la definición de dispositivo "
        "médico: la pregunta aplica.\n\n"
        "No invasivo · piel intacta · ilumina en espectro visible\n"
        "→ Clase I, registro simplificado (semanas, no meses).\n\n"
        "Con luz ultravioleta la ruta sería mucho más larga.",
        CIAN, tam_c=12)

    _tarjeta_texto(
        s, Inches(6.95), Inches(3.85), Inches(5.6), Inches(2.75),
        "RIESGOS Y MITIGACIÓN",
        "Tranquilizar de más → IMPOSIBLE POR DISEÑO: la fiebre escala y "
        "nunca rebaja.\n\n"
        "Esa regla baja la clase de seguridad del software de C a B "
        "(IEC 62304).\n\n"
        "Al revisar ISO 14971 los controles que exige ya estaban en el código.",
        MAGENTA, tam_c=12)
    _notas(s, """
[3:10–4:00] VIABILIDAD — 50 s

«El costo por unidad es de 226 a 401 soles. Un citómetro de flujo cuesta decenas
de miles de dólares. Y no consume reactivos: el costo de operación después de la
primera unidad es cero. Funciona sin internet, porque una app que necesita
conexión no sirve en Amazonas.

Y sobre si esto se puede registrar: la Ley 29459 incluye al software en la
definición de dispositivo médico, así que la pregunta aplica. Por los criterios
de DIGEMID somos CLASE I — no invasivo, piel intacta, e ilumina en el espectro
visible, que es una excepción explícita de la norma. Registro simplificado.
Si hubiéramos elegido ultravioleta, la ruta sería mucho más larga.

Y algo que nos sorprendió: al revisar la ISO 14971 resultó que las decisiones de
seguridad que ya habíamos tomado son exactamente los controles que la norma
exige. Que la fiebre nunca pueda rebajar la conducta no es solo prudencia
clínica: baja la clase de seguridad del software de C a B.»
""")

    s = prs.slides.add_slide(blanco); _fondo(s)
    _encabezado(s, "6.", "Enfoque en el usuario y alternativas", VIOLETA)
    _texto(s, Inches(.78), Inches(1.4), Inches(11.8), Inches(.45),
           "Dos aplicaciones, un solo sistema — porque son dos usuarios "
           "con dos necesidades distintas",
           tam=15, color=VIOLETA, negrita=True)

    if (CAP / "clinica_ventana.jpg").exists():
        _figura(s, CAP / "clinica_ventana.jpg", Inches(1.0), Inches(2.35),
                Inches(5.1), h_max=Inches(4.15))
    if (CAP / "familia_racha.jpg").exists():
        _figura(s, CAP / "familia_racha.jpg", Inches(7.0), Inches(2.35),
                Inches(5.1), h_max=Inches(4.15))

    _texto(s, Inches(1.0), Inches(1.95), Inches(5.1), Inches(.4),
           "LA POSTA · enfermera o técnico", tam=12, color=CIAN, negrita=True)
    _texto(s, Inches(7.0), Inches(1.95), Inches(5.1), Inches(.4),
           "LA FAMILIA · todos los días", tam=12, color=MAGENTA, negrita=True)
    _pie(s, "Contexto real: sin laboratorio, sin hematólogo, con señal "
            "intermitente · 28 KB, funciona offline")
    _notas(s, """
[4:00–4:50] USUARIO — 50 s   ← DEMO EN VIVO, no capturas

PARTE A · la posta (15 s):
Escribir edad 6 → APARECE SOLO el aviso pediátrico. Analizar → semáforo NEGRO.
«Fíjense: el ajuste pediátrico aparece solo, sin que nadie lo pida. Y no entrega
un número, entrega una conducta — con el plazo, y diciendo que el antibiótico va
ANTES del traslado, no al llegar, porque el sistema sabe que el centro está a
nueve horas.»

PARTE B · la familia (25 s) — EN UN CELULAR REAL:
Pulsar «Ya se la dimos». Confeti, sonido, vibración, racha.
«Y esta la usan los papás, todos los días. Es la pastilla de la noche, con las
tres reglas que casi nadie conoce: DE NOCHE Y NO DE MAÑANA, SIN LÁCTEOS, Y
SIEMPRE IGUAL. Las tres tienen evidencia detrás y las tres cambian si el
medicamento funciona.»

⚠️ LLEVAR VÍDEO DE PANTALLA GRABADO COMO RESPALDO.
""")

    s = prs.slides.add_slide(blanco); _fondo(s)
    _encabezado(s, "7.", "Roadmap, escalabilidad y adaptabilidad", NARANJA)

    for i, (plazo, texto, color) in enumerate([
        ("CORTO · 1–6 meses",
         "Comité de ética del INSNSB\n+ 30–50 vídeos con hemograma pareado",
         MAGENTA),
        ("MEDIANO · 6–12 meses",
         "Umbrales pediátricos definitivos\n+ piloto en 2–3 postas",
         VIOLETA),
        ("LARGO · 12–24 meses",
         "Registro DIGEMID\n+ integración a RENHICE",
         CIAN),
    ]):
        x = Inches(.8 + i * 4.05)
        _tarjeta(s, x, Inches(1.6), Inches(3.75), Inches(1.95), color)
        _texto(s, x + Inches(.25), Inches(1.8), Inches(3.3), Inches(.4),
               plazo, tam=12, color=color, negrita=True)
        _texto(s, x + Inches(.25), Inches(2.28), Inches(3.3), Inches(1.1),
               texto, tam=13, interlineado=1.25)
        if i < 2:
            p = slide_punto = s.shapes.add_shape(
                MSO_SHAPE.OVAL, x + Inches(3.85), Inches(2.45),
                Inches(.14), Inches(.14))
            p.fill.solid(); p.fill.fore_color.rgb = GRIS
            p.line.fill.background(); p.shadow.inherit = False

    _figura(s, FIG / "08_tres_senales.png", Inches(1.35), Inches(3.9),
            Inches(7.1), h_max=Inches(2.6))
    _tarjeta_texto(
        s, Inches(8.95), Inches(3.78), Inches(3.6), Inches(2.6),
        "ESCALABILIDAD",
        "El mismo equipo sirve a CUALQUIER paciente en quimioterapia, "
        "no solo LLA.\n\n"
        "Y del mismo vídeo salen más señales, sin hardware adicional.",
        VERDE, tam_c=12)
    _notas(s, """
[4:50–5:20] ROADMAP — 30 s

«Lo que falta no es más código. Son 30 o 50 VÍDEOS REALES CON HEMOGRAMA PAREADO.
Ese es el pedido que le hacemos al instituto.

Y escala solo: el mismo aparato sirve para cualquier niño en quimioterapia.
Del mismo vídeo salen más señales — hemoglobina con un segundo LED, y perfusión
microcirculatoria sin hardware adicional.»

NO profundizar en microcirculación aquí. Si preguntan, está la reserva.
""")

    s = prs.slides.add_slide(blanco); _fondo(s)
    _encabezado(s, "8.", "Equipo", CIAN)
    _texto(s, Inches(.78), Inches(1.4), Inches(11.8), Inches(.5),
           "Medicina · Bioingeniería · Mecatrónica · Computación",
           tam=17, color=GRIS)

    roles = [
        ("Liderazgo", "[nombre]", "Coordina y define prioridades", MAGENTA),
        ("Clínico", "[nombre]", "Umbrales, criterios de derivación", CIAN),
        ("Tecnología", "[nombre]", "Pipeline, modelo, apps", VIOLETA),
        ("Diseño / UX", "[nombre]", "Las dos interfaces y el michi", NARANJA),
        ("Investigación", "[nombre]", "Óptica, hardware, evidencia", VERDE),
    ]
    for i, (rol, nombre, aporta, color) in enumerate(roles):
        x = Inches(.8 + i * 2.42)
        _tarjeta(s, x, Inches(2.15), Inches(2.2), Inches(2.35), color)
        _texto(s, x + Inches(.2), Inches(2.38), Inches(1.85), Inches(.4),
               rol.upper(), tam=11, color=color, negrita=True)
        _texto(s, x + Inches(.2), Inches(2.8), Inches(1.85), Inches(.5),
               nombre, tam=14, negrita=True)
        _texto(s, x + Inches(.2), Inches(3.35), Inches(1.85), Inches(1.0),
               aporta, tam=11, color=GRIS, interlineado=1.2)

    _tarjeta(s, Inches(.8), Inches(4.85), Inches(11.75), Inches(1.6), NARANJA)
    _texto(s, Inches(1.2), Inches(5.08), Inches(11.0), Inches(1.2),
           "Lo que hicimos estos días fue, sobre todo, medir en vez de suponer:\n"
           "que a 30 fps el método no funciona · que el umbral del adulto se "
           "pierde la franja crítica · que nuestro modelo entrenado era peor "
           "que la física sola, así que lo cambiamos.",
           tam=15, interlineado=1.35, alinear=PP_ALIGN.CENTER)
    _notas(s, """
[5:20] CIERRE — 20 s

«Somos un equipo de medicina, bioingeniería, mecatrónica y computación. Y lo que
hicimos estos días fue, sobre todo, MEDIR EN VEZ DE SUPONER.

La niña de Bagua no necesita un hospital más cerca. Necesita que alguien en su
posta pueda responder una pregunta. Eso es MichiCheck.»
""")

    reservas = [
        ("04_metricas.png", "Precisión del modelo",
         "AUC 0.939 · Sensibilidad 94.7 % · VPN 95.1 %\n"
         "Validación cruzada por paciente, n=300, cohorte sintética.\n"
         "Mide coherencia del pipeline, NO exactitud clínica.",
         "«¿Qué precisión tiene?»"),
        ("05_validacion_real.png", "Validación con vídeo real",
         "Estabilización, segmentación y ajuste de diámetro validados sobre "
         "vídeo capilaroscópico público que no generamos nosotros.\n"
         "La cadena completa hasta el ANC, no: falta escala conocida y "
         "hemograma pareado.",
         "«¿Lo probaron con datos reales?»"),
        ("03_requisito_fps.png", "Por qué ≥60 fps",
         "A 30 fps el error de velocimetría es del 81 %: el método no funciona.\n"
         "A 60 fps, 7 %. El requisito no es resolución, es tasa de fotogramas.",
         "«¿Por qué necesitan esa cámara?»"),
        ("06_ventana_terapeutica.png", "Ventana terapéutica y adherencia",
         "En mantenimiento el objetivo ES 500–1500. Un recuento alto significa "
         "que el tratamiento no está llegando.\n"
         "44 % de adherencia insuficiente × 2.7 de riesgo de recaída.",
         "«¿Cómo detectan la no adherencia?»"),
    ]
    for archivo, titulo, cuerpo, cuando in reservas:
        ruta = FIG / archivo
        if not ruta.exists():
            continue
        s = prs.slides.add_slide(blanco); _fondo(s)
        _texto(s, Inches(.78), Inches(.4), Inches(11.8), Inches(.5),
               "RESERVA", tam=11, color=NARANJA, negrita=True)
        _texto(s, Inches(.78), Inches(.75), Inches(11.8), Inches(.6),
               titulo, tam=28, negrita=True)
        _figura(s, ruta, Inches(1.4), Inches(1.75), Inches(10.5),
                h_max=Inches(4.3))
        _texto(s, Inches(.9), Inches(6.25), Inches(11.5), Inches(1.0),
               cuerpo, tam=13, color=GRIS, alinear=PP_ALIGN.CENTER,
               interlineado=1.25)
        _notas(s, f"Diapositiva de reserva. Sacarla si preguntan: {cuando}")

    return prs


if __name__ == "__main__":
    prs = construir()
    prs.save(SALIDA)
    print(f"{SALIDA}  ({SALIDA.stat().st_size / 1e6:.1f} MB, "
          f"{len(prs.slides._sldIdLst)} diapositivas)")
